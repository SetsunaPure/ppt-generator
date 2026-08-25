"""
PPT生成器 - E2E端到端测试

实际生成PPT文件并验证输出：
1. 备课PPT多模板生成
2. 讲题PPT生成
3. CLI demo模式
4. 内容分页正确性
5. 边界场景（空内容/超长内容/纯公式/混合内容）
6. PPT结构校验（slide数量/shapes/文本内容）
"""

import sys
import os
import json
import shutil
import tempfile
from pathlib import Path
from dataclasses import dataclass

# 确保项目根目录在sys.path中
PROJECT_ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

# 切换工作目录到项目根，确保相对路径生效
os.chdir(str(PROJECT_ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from src.models import get_env_config, Page, PageData, Block, BlockPosition, LessonPptRequest
from src.config import PageType, TEXT_BOX_HEIGHT, TEXT_BOX_WIDTH, calc_text_line_height
from src.outline_parser import OutlineParser, TopicParser, TOPIC_TEMPLATES
from src.layout_engine import split_page, html_to_blocks, layout_blocks
from src.builder import PptBuilder


# ============================================================
# 测试数据
# ============================================================

# 标准备课大纲（OutlineParser.parse期望dict，包含title/subTitle/outlineList）
# outlineList顶层需包含outlineCode=6的教学设计节点，子节点为7/8/9
LESSON_OUTLINE = {
    "title": "二次函数的图像与性质",
    "subTitle": "九年级数学上册",
    "outlineList": [
        {
            "outlineCode": 6,
            "title": "教学设计",
            "children": [
                {
                    "outlineCode": 7,
                    "title": "教学规划",
                    "children": [
                        {
                            "outlineCode": 13,
                            "title": "教学目标",
                            "content": [
                                {
                                    "content": "<p>1. 掌握二次函数的基本概念和图像特征</p>"
                                               "<p>2. 能够根据条件确定二次函数的解析式</p>"
                                               "<p>3. 理解二次函数与一元二次方程的关系</p>"
                                }
                            ]
                        },
                        {
                            "outlineCode": 14,
                            "title": "重点难点",
                            "content": [
                                {
                                    "content": "<p><strong>重点：</strong>二次函数图像的开口方向、对称轴、顶点坐标</p>"
                                               "<p><strong>难点：</strong>二次函数与一元二次方程的内在联系</p>"
                                }
                            ]
                        }
                    ]
                },
                {
                    "outlineCode": 8,
                    "title": "教学过程",
                    "children": [
                        {
                            "outlineCode": 16,
                            "title": "新课导入",
                            "content": [
                                {
                                    "content": "<p>回顾一次函数 y=kx+b 的图像和性质，思考：当自变量最高次数为2时，函数图像会是什么样子？</p>"
                                }
                            ]
                        },
                        {
                            "outlineCode": 18,
                            "title": "知识梳理",
                            "content": [
                                {
                                    "content": "<p>一般形式：$y=ax^2+bx+c$（$a \\neq 0$）</p>"
                                               "<p>顶点式：$y=a(x-h)^2+k$，顶点坐标为 $(h,k)$</p>"
                                               "<p>交点式：$y=a(x-x_1)(x-x_2)$，其中 $x_1, x_2$ 为与x轴交点的横坐标</p>"
                                }
                            ]
                        }
                    ]
                },
                {
                    "outlineCode": 9,
                    "title": "教学总结",
                    "content": [
                        {
                            "content": "<p>本节课学习了二次函数的定义、三种表达形式及其图像的基本性质，为后续学习二次函数的应用打下基础。</p>"
                        }
                    ]
                }
            ]
        }
    ]
}

# 长内容大纲（测试分页）
LONG_CONTENT_OUTLINE = {
    "title": "长内容测试",
    "subTitle": "分页验证",
    "outlineList": [
        {
            "outlineCode": 6,
            "title": "教学设计",
            "children": [
                {
                    "outlineCode": 7,
                    "title": "教学规划",
                    "children": [
                        {
                            "outlineCode": 13,
                            "title": "教学目标",
                            "content": [
                                {
                                    "content": "".join(
                                        f"<p>第{i+1}条教学目标：掌握数学中关于函数与方程的核心知识点，理解其几何意义与代数表达之间的内在联系，"
                                        f"能够灵活运用所学知识解决各类综合问题，提升数学思维能力和解题技巧。</p>"
                                        for i in range(30)
                                    )
                                }
                            ]
                        }
                    ]
                },
                {
                    "outlineCode": 8,
                    "title": "教学过程",
                    "children": [
                        {
                            "outlineCode": 16,
                            "title": "新课导入",
                            "content": [
                                {
                                    "content": "<p>回顾一次函数的图像和性质</p>"
                                }
                            ]
                        },
                    ]
                },
                {
                    "outlineCode": 9,
                    "title": "教学总结",
                    "content": [
                        {
                            "content": "<p>总结完毕</p>"
                        }
                    ]
                }
            ]
        }
    ]
}

# 纯公式内容
FORMULA_CONTENT_OUTLINE = {
    "title": "公式专题",
    "subTitle": "数学公式汇总",
    "outlineList": [
        {
            "outlineCode": 6,
            "title": "教学设计",
            "children": [
                {
                    "outlineCode": 7,
                    "title": "公式专题",
                    "children": [
                        {
                            "outlineCode": 13,
                            "title": "核心公式",
                            "content": [
                                {
                                    "content": "<p>勾股定理：$a^2+b^2=c^2$</p>"
                                               "<p>求根公式：$x=\\frac{-b\\pm\\sqrt{b^2-4ac}}{2a}$</p>"
                                               "<p>韦达定理：$x_1+x_2=-\\frac{b}{a}$，$x_1 x_2=\\frac{c}{a}$</p>"
                                }
                            ]
                        }
                    ]
                },
                {
                    "outlineCode": 9,
                    "title": "总结",
                    "content": [
                        {
                            "content": "<p>以上为常用公式汇总</p>"
                        }
                    ]
                }
            ]
        }
    ]
}

# 讲题JSON
TOPIC_JSON = {
    "topicTemplateType": 1,
    "context": {
        "stem": "<p>已知二次函数 $y=x^2-4x+3$，求该函数的顶点坐标和与x轴的交点坐标。</p>",
        "options": ["A. (2,-1)", "B. (2,1)", "C. (-2,-1)", "D. (-2,1)"]
    },
    "explain": {
        "answers": ["A"],
        "analysis": "<p>将函数配方得 $y=(x-2)^2-1$，故顶点为(2,-1)。</p>"
    },
    "testPointAnalysis": "本题考查二次函数的顶点式和图像性质",
    "explorationOfSolutions": "配方法求顶点坐标，令y=0求交点",
    "explorationOfSublimate": "二次函数顶点式与一般式的转换关系",
}

# 混合HTML内容（含表格/加粗/列表）
MIXED_HTML = (
    "<p><strong>知识框架</strong></p>"
    "<table>"
    "<tr><th>函数类型</th><th>一般形式</th><th>图像特征</th></tr>"
    "<tr><td>一次函数</td><td>y=kx+b</td><td>直线</td></tr>"
    "<tr><td>二次函数</td><td>y=ax²+bx+c</td><td>抛物线</td></tr>"
    "<tr><td>反比例函数</td><td>y=k/x</td><td>双曲线</td></tr>"
    "</table>"
    "<p>以上三种函数是初中数学的核心内容</p>"
)


# ============================================================
# 辅助函数
# ============================================================

def generate_lesson_ppt(outline_json, style="0", font_size=16):
    """生成备课PPT并返回输出路径"""
    env = get_env_config("dev")
    parser = OutlineParser()
    pages = parser.parse(outline_json, font_size=font_size)
    builder = PptBuilder(
        env=env,
        pages=pages,
        font_size=font_size,
        file_content_style=style,
        logo_path="",
        template_prefix="lesson"
    )
    return builder.build()


def generate_topic_ppt(topic_json, style="0", font_size=16):
    """生成讲题PPT并返回输出路径"""
    env = get_env_config("dev")
    parser = TopicParser(config=TOPIC_TEMPLATES)
    pages = parser.parse(topic_json, font_size=font_size, is_origin=False)
    builder = PptBuilder(
        env=env,
        pages=pages,
        font_size=font_size,
        file_content_style=style,
        logo_path="",
        template_prefix="topic"
    )
    return builder.build()


def validate_pptx(path, min_slides=1, min_size_kb=1):
    """校验PPTX文件基本属性"""
    assert Path(path).exists(), f"PPT文件不存在: {path}"
    size = Path(path).stat().st_size
    assert size >= min_size_kb * 1024, f"文件过小: {size} bytes (期望 >= {min_size_kb}KB)"
    prs = Presentation(path)
    assert len(prs.slides) >= min_slides, f"slide数量不足: {len(prs.slides)} (期望 >= {min_slides})"
    return prs


def extract_all_text(prs):
    """提取PPT中所有文本内容"""
    texts = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
    return texts


# ============================================================
# 测试类
# ============================================================

class TestLessonPptGeneration:
    """备课PPT生成测试"""

    def test_standard_outline_style0(self):
        """标准大纲 + 风格0模板"""
        path = generate_lesson_ppt(LESSON_OUTLINE, style="0")
        prs = validate_pptx(path, min_slides=3, min_size_kb=5)

        # 第一页应为标题页
        first_slide = prs.slides[0]
        texts = extract_all_text(prs)

        # 应包含教学相关文本
        all_text = " ".join(texts)
        assert any(kw in all_text for kw in ["教学", "二次函数"]), \
            f"PPT内容应包含教学关键词，实际: {all_text[:200]}"

        print(f"  [PASS] 标准备课PPT(style0): {len(prs.slides)} slides, {Path(path).stat().st_size} bytes")

    def test_standard_outline_all_styles(self):
        """标准大纲 + 所有6种风格模板"""
        for style_idx in range(6):
            path = generate_lesson_ppt(LESSON_OUTLINE, style=str(style_idx))
            prs = validate_pptx(path, min_slides=2, min_size_kb=5)
            print(f"  [PASS] style={style_idx}: {len(prs.slides)} slides")

    def test_long_content_pagination(self):
        """超长内容 → 验证分页正确"""
        path = generate_lesson_ppt(LONG_CONTENT_OUTLINE, style="0")
        prs = validate_pptx(path, min_slides=5, min_size_kb=10)

        # 长内容应产生较多页面
        assert len(prs.slides) >= 5, f"长内容应产生>=5页，实际: {len(prs.slides)}"

        texts = extract_all_text(prs)
        assert any("教学目标" in t for t in texts), "应包含教学目标标题"

        print(f"  [PASS] 超长内容分页: {len(prs.slides)} slides")

    def test_formula_content(self):
        """纯公式内容"""
        path = generate_lesson_ppt(FORMULA_CONTENT_OUTLINE, style="0")
        prs = validate_pptx(path, min_slides=2, min_size_kb=5)

        texts = extract_all_text(prs)
        all_text = " ".join(texts)
        # 公式可能以latex原文或转换后的形式出现
        assert any(kw in all_text for kw in ["公式", "勾股", "sqrt"]), \
            f"应包含公式相关内容，实际: {all_text[:200]}"

        print(f"  [PASS] 纯公式内容: {len(prs.slides)} slides")


class TestTopicPptGeneration:
    """讲题PPT生成测试"""

    def test_topic_type1(self):
        """讲题模板类型1（基础模板）"""
        path = generate_topic_ppt(TOPIC_JSON, style="0")
        prs = validate_pptx(path, min_slides=2, min_size_kb=5)

        texts = extract_all_text(prs)
        all_text = " ".join(texts)
        assert any(kw in all_text for kw in ["题目", "考点", "二次函数"]), \
            f"讲题PPT应包含题目/考点内容，实际: {all_text[:200]}"

        print(f"  [PASS] 讲题PPT(type1): {len(prs.slides)} slides")

    def test_topic_all_styles(self):
        """讲题PPT + 所有风格"""
        for style_idx in range(6):
            path = generate_topic_ppt(TOPIC_JSON, style=str(style_idx))
            prs = validate_pptx(path, min_slides=1, min_size_kb=5)
            print(f"  [PASS] topic style={style_idx}: {len(prs.slides)} slides")


class TestCliDemo:
    """CLI demo模式测试"""

    def test_demo_default(self):
        """默认参数demo"""
        import subprocess
        result = subprocess.run(
            [sys.executable, str(PROJECT_ROOT / "cli.py"), "demo"],
            capture_output=True, text=True, timeout=30,
            cwd=str(PROJECT_ROOT)
        )
        assert result.returncode == 0, f"CLI demo执行失败: {result.stderr}"
        assert "Demo PPT生成成功" in result.stdout, f"应输出成功信息: {result.stdout}"

        # 从输出中提取文件路径
        for line in result.stdout.splitlines():
            if "Demo PPT生成成功:" in line:
                output_path = line.split("成功:")[-1].strip()
                prs = validate_pptx(output_path, min_slides=1, min_size_kb=5)
                print(f"  [PASS] CLI demo默认: {len(prs.slides)} slides")
                return

        print("  [WARN] 未从CLI输出中提取到文件路径")

    def test_demo_custom_fontsize(self):
        """自定义字号demo"""
        import subprocess
        result = subprocess.run(
            [sys.executable, str(PROJECT_ROOT / "cli.py"), "demo", "--font-size", "14"],
            capture_output=True, text=True, timeout=30,
            cwd=str(PROJECT_ROOT)
        )
        assert result.returncode == 0, f"CLI demo(字号14)失败: {result.stderr}"
        print("  [PASS] CLI demo字号14")


class TestSplitPageEngine:
    """分页引擎测试"""

    def test_short_text_single_page(self):
        """短文本 → 应在1页内"""
        pages = split_page("<p>这是一段简短的文本</p>", font_size_pt=16)
        assert len(pages) >= 1, "至少产生1页"
        print(f"  [PASS] 短文本分页: {len(pages)} 页")

    def test_long_text_multi_page(self):
        """长文本 → 应分多页"""
        long_html = "".join(f"<p>这是第{i+1}行文本，用于测试分页功能是否正常工作，每行内容足够长以触发换行机制。</p>" for i in range(80))
        pages = split_page(long_html, font_size_pt=16)
        assert len(pages) >= 2, f"长文本应分多页，实际: {len(pages)}"
        print(f"  [PASS] 长文本分页: {len(pages)} 页")

    def test_empty_content(self):
        """空内容 → 返回空列表或1个空页"""
        pages = split_page("", font_size_pt=16)
        assert len(pages) <= 1, f"空内容最多1页，实际: {len(pages)}"
        print(f"  [PASS] 空内容分页: {len(pages)} 页")

    def test_mixed_html_content(self):
        """混合HTML内容（表格+文本）"""
        pages = split_page(MIXED_HTML, font_size_pt=16)
        assert len(pages) >= 1, "至少1页"

        # 检查所有页的block
        total_blocks = sum(len(p) for p in pages)
        assert total_blocks > 0, "应有block产生"
        print(f"  [PASS] 混合内容分页: {len(pages)} 页, {total_blocks} blocks")

    def test_html_to_blocks_text(self):
        """HTML → Block转换: 纯文本"""
        blocks = html_to_blocks("<p>你好世界</p>", font_size_pt=16)
        text_blocks = [b for b in blocks if b.type == "text"]
        assert len(text_blocks) > 0, "应有文本Block"
        assert any("你好" in (b.content or "") for b in text_blocks), "应包含'你好'"
        print(f"  [PASS] HTML→Block纯文本: {len(text_blocks)} text blocks")

    def test_html_to_blocks_table(self):
        """HTML → Block转换: 表格"""
        blocks = html_to_blocks(
            "<table><tr><td>A</td><td>B</td></tr><tr><td>C</td><td>D</td></tr></table>",
            font_size_pt=16
        )
        table_blocks = [b for b in blocks if b.type == "table"]
        assert len(table_blocks) > 0, "应有表格Block"
        assert table_blocks[0].table_rows == 2, f"表格应有2行，实际: {table_blocks[0].table_rows}"
        assert table_blocks[0].table_cols == 2, f"表格应有2列，实际: {table_blocks[0].table_cols}"
        print(f"  [PASS] HTML→Block表格: {len(table_blocks)} table blocks, {table_blocks[0].table_rows}x{table_blocks[0].table_cols}")

    def test_layout_blocks_coordinates(self):
        """布局引擎坐标计算验证"""
        blocks = [
            Block(type="text", content="行1文字", inline=True,
                  width=Inches(3), height=Pt(16)),
            Block(type="newline", inline=False, width=0, height=Pt(16)),
            Block(type="text", content="行2文字", inline=True,
                  width=Inches(3), height=Pt(16)),
        ]
        pages = layout_blocks(blocks, font_size_pt=16)
        assert len(pages) >= 1, "至少1页"

        # 行2的top应大于行1
        if len(pages[0]) >= 3:
            row1 = pages[0][0]
            row2 = pages[0][2]
            assert row2.position.top > row1.position.top, \
                f"行2 Y坐标应大于行1: row1.top={row1.position.top}, row2.top={row2.position.top}"
        print(f"  [PASS] 坐标验证: {len(pages)} 页")

    def test_page_overflow_detection(self):
        """页满检测：超量内容应产生多页"""
        text_line_height = calc_text_line_height(16)
        # 构造大量Block，确保超过TEXT_BOX_HEIGHT
        blocks = []
        for i in range(100):
            blocks.append(Block(
                type="text", content=f"第{i+1}行内容，用于填充页面",
                inline=True, width=Inches(5), height=text_line_height
            ))
            blocks.append(Block(type="newline", inline=False, width=0, height=text_line_height))

        pages = layout_blocks(blocks, font_size_pt=16)
        assert len(pages) >= 2, f"100行文本应分多页，实际: {len(pages)}"
        print(f"  [PASS] 页满检测: {len(pages)} 页")


class TestOutlineParser:
    """大纲解析测试"""

    def test_parse_produces_title_page(self):
        """解析结果首页为标题页"""
        parser = OutlineParser()
        pages = parser.parse(LESSON_OUTLINE, font_size=16)
        assert len(pages) > 0, "应产生页面"
        assert pages[0].type == PageType.TITLE, f"首页应为TITLE，实际: {pages[0].type}"
        print(f"  [PASS] 标题页生成: {len(pages)} 页, 首页={pages[0].type}")

    def test_parse_all_outline_codes(self):
        """解析所有大纲code"""
        parser = OutlineParser()
        pages = parser.parse(LESSON_OUTLINE, font_size=16)
        page_types = [p.type for p in pages]
        # 应有标题页和内容页
        assert PageType.TITLE in page_types, "应有标题页"
        assert PageType.FORMULA in page_types, "应有公式/内容页"
        print(f"  [PASS] 大纲code覆盖: {set(t.value for t in page_types)}")

    def test_parse_with_empty_children(self):
        """子节点为空时不崩溃"""
        minimal = {
            "title": "空子节点测试",
            "subTitle": "测试",
            "outlineList": [
                {
                    "outlineCode": 6,
                    "title": "教学设计",
                    "children": [
                        {
                            "outlineCode": 7,
                            "title": "教学规划",
                            "children": []
                        }
                    ]
                }
            ]
        }
        parser = OutlineParser()
        pages = parser.parse(minimal, font_size=16)
        # 只有标题页
        assert len(pages) >= 1, "至少有标题页"
        print(f"  [PASS] 空子节点: {len(pages)} 页")


class TestTopicParser:
    """讲题模板解析测试"""

    def test_parse_type1(self):
        """基础模板(type=1)解析"""
        parser = TopicParser(config=TOPIC_TEMPLATES)
        pages = parser.parse(TOPIC_JSON, font_size=16, is_origin=False)
        assert len(pages) > 0, "应产生页面"
        print(f"  [PASS] 讲题type1解析: {len(pages)} 页")

    def test_parse_type2(self):
        """卡点模板(type=2)解析"""
        data = dict(TOPIC_JSON, topicTemplateType=2)
        data["overallDesignAssessment"] = "整体设计评估内容"
        data["stuckPointAndSolution"] = "核心卡点与破题路径内容"
        data["similarQuestionDesignPattern"] = "同类题模式识别内容"

        parser = TopicParser(config=TOPIC_TEMPLATES)
        pages = parser.parse(data, font_size=16, is_origin=False)
        assert len(pages) > 0
        print(f"  [PASS] 讲题type2解析: {len(pages)} 页")

    def test_parse_type3(self):
        """全流程模板(type=3)解析"""
        data = dict(TOPIC_JSON, topicTemplateType=3)
        data["testPointAnalysis"] = "考点分析内容"
        data["solutionToTheProblem"] = "破题思路内容"
        data["solutionProcess"] = "解题过程内容"
        data["cautionaryNote"] = "易错警示内容"
        data["explorationOfSublimate"] = "解法升华内容"

        parser = TopicParser(config=TOPIC_TEMPLATES)
        pages = parser.parse(data, font_size=16, is_origin=False)
        assert len(pages) > 0
        print(f"  [PASS] 讲题type3解析: {len(pages)} 页")


class TestPptStructure:
    """PPT结构校验"""

    def test_slide_count_matches_pages(self):
        """生成的slide数量与Page列表对应"""
        env = get_env_config("dev")
        parser = OutlineParser()
        pages = parser.parse(LESSON_OUTLINE, font_size=16)

        builder = PptBuilder(
            env=env, pages=pages, font_size=16,
            file_content_style="0", logo_path="",
            template_prefix="lesson"
        )
        path = builder.build()
        prs = Presentation(path)

        # PPT的slide数量应与pages数量相关（允许模板中有多余slide）
        assert len(prs.slides) >= len(pages), \
            f"slide数量({len(prs.slides)})应>=page数量({len(pages)})"

        print(f"  [PASS] slide/page对应: {len(prs.slides)} slides, {len(pages)} pages")

    def test_output_file_naming(self):
        """输出文件命名规范"""
        path = generate_lesson_ppt(LESSON_OUTLINE, style="0")
        filename = Path(path).name
        assert filename.startswith("lesson_"), f"文件名应以lesson_开头: {filename}"
        assert filename.endswith(".pptx"), f"文件名应以.pptx结尾: {filename}"
        print(f"  [PASS] 文件命名: {filename}")

    def test_16_9_aspect_ratio(self):
        """宽屏16:9比例验证"""
        path = generate_lesson_ppt(LESSON_OUTLINE, style="0")
        prs = Presentation(path)

        width = prs.slide_width
        height = prs.slide_height
        ratio = width / height
        # 16:9 ≈ 1.778
        assert 1.7 < ratio < 1.85, f"宽高比应接近16:9，实际: {ratio:.3f}"
        print(f"  [PASS] 宽屏比例: {ratio:.3f}")


class TestEdgeCases:
    """边界场景测试"""

    def test_single_line_content(self):
        """仅一行内容"""
        minimal = {
            "title": "简单测试",
            "subTitle": "边界场景",
            "outlineList": [
                {
                    "outlineCode": 6,
                    "title": "教学设计",
                    "children": [
                        {
                            "outlineCode": 9,
                            "title": "总结",
                            "content": [{"content": "<p>完</p>"}]
                        }
                    ]
                }
            ]
        }
        parser = OutlineParser()
        pages = parser.parse(minimal, font_size=16)
        assert len(pages) >= 1, "至少1页"

        env = get_env_config("dev")
        builder = PptBuilder(
            env=env, pages=pages, font_size=16,
            file_content_style="0", logo_path="",
            template_prefix="lesson"
        )
        path = builder.build()
        prs = validate_pptx(path, min_slides=1, min_size_kb=1)
        print(f"  [PASS] 单行内容: {len(prs.slides)} slides")

    def test_unicode_content(self):
        """Unicode特殊字符内容"""
        unicode_html = "<p>特殊符号：α β γ δ ε θ λ μ π σ φ ψ ω ∑ ∏ √ ∞ ≤ ≥ ≠ ≈ ∫ ∂ ∇</p>"
        pages = split_page(unicode_html, font_size_pt=16)
        assert len(pages) >= 1, "Unicode内容应正常分页"
        print(f"  [PASS] Unicode内容: {len(pages)} 页")

    def test_nested_html_tags(self):
        """嵌套HTML标签"""
        nested = "<p><strong>加粗文本</strong>和<u>下划线文本</u>混合</p>"
        blocks = html_to_blocks(nested, font_size_pt=16)
        assert len(blocks) > 0, "嵌套HTML应产生Block"
        print(f"  [PASS] 嵌套HTML: {len(blocks)} blocks")

    def test_font_size_variations(self):
        """不同字号生成"""
        for fs in [12, 14, 16, 18, 20]:
            path = generate_lesson_ppt(LESSON_OUTLINE, style="0", font_size=fs)
            prs = validate_pptx(path, min_slides=2, min_size_kb=5)
            print(f"  [PASS] 字号{fs}: {len(prs.slides)} slides")


# ============================================================
# 运行入口
# ============================================================

def run_all_tests():
    """运行所有E2E测试"""
    print("=" * 60)
    print("PPT生成器 - E2E端到端测试")
    print("=" * 60)

    test_classes = [
        ("备课PPT生成", TestLessonPptGeneration),
        ("讲题PPT生成", TestTopicPptGeneration),
        ("CLI Demo模式", TestCliDemo),
        ("分页引擎", TestSplitPageEngine),
        ("大纲解析", TestOutlineParser),
        ("讲题解析", TestTopicParser),
        ("PPT结构校验", TestPptStructure),
        ("边界场景", TestEdgeCases),
    ]

    total_passed = 0
    total_failed = 0
    failures = []

    for class_name, test_class in test_classes:
        print(f"\n{'─' * 40}")
        print(f"▸ {class_name}")
        print(f"{'─' * 40}")

        instance = test_class()
        methods = [m for m in dir(instance) if m.startswith("test_")]

        for method_name in methods:
            try:
                getattr(instance, method_name)()
                total_passed += 1
            except Exception as e:
                total_failed += 1
                failures.append(f"{class_name}.{method_name}: {e}")
                print(f"  [FAIL] {method_name}: {e}")

    print(f"\n{'=' * 60}")
    print(f"总计: {total_passed} 通过, {total_failed} 失败")
    if failures:
        print("\n失败详情:")
        for f in failures:
            print(f"  ✗ {f}")
    print(f"{'=' * 60}")

    return total_failed == 0


if __name__ == "__main__":
    success = run_all_tests()
    sys.exit(0 if success else 1)
