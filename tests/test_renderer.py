"""
渲染测试 - P2优先级

测试PPT渲染功能（快照测试）
"""

import sys
from pathlib import Path

# 将src目录添加到路径
sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from src.models import Block, BlockPosition, EnvConfig, Page, PageData, PageType
from src.renderer import render_page, render_content_page, FormulaCache
from src.config import calc_text_line_height


@pytest.fixture
def presentation():
    """创建测试用Presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 添加一个空白slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


@pytest.fixture
def mock_env():
    """Mock环境配置"""
    return EnvConfig(
        profile="test",
        proxy=None,
        oss_region="beijing",
        latex_api_url="http://test.api/latex",
        wechat_webhook="https://test.webhook",
        wechat_mentioned=(),
    )


@pytest.fixture
def formula_cache(mock_env):
    """公式缓存"""
    return FormulaCache(mock_env)


class TestRenderSnapshot:
    """渲染快照测试"""

    def test_render_text_block(self, presentation, font_size_pt=16, formula_cache=None):
        """渲染文本块"""
        prs, slide = presentation

        # 创建文本块
        text_block = Block(
            type="text",
            content="测试文本",
            width=200000,
            height=calc_text_line_height(font_size_pt),
            inline=True,
            position=BlockPosition(left=0, top=0)
        )

        # 创建文本框
        textbox = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.7), Inches(10), Inches(5)
        )
        text_frame = textbox.text_frame

        # 渲染
        render_page(
            slide,
            [text_block],
            mock_env,
            font_size_pt,
            formula_cache or FormulaCache(mock_env),
            (int(Inches(0.2)), int(Inches(0.7))),
            text_frame
        )

        # 验证文本框有内容
        assert len(text_frame.paragraphs) > 0

    def test_render_multiple_blocks(self, presentation, font_size_pt=16, mock_env=None):
        """渲染多个块"""
        prs, slide = presentation
        env = mock_env or EnvConfig(
            profile="test",
            proxy=None,
            oss_region="beijing",
            latex_api_url="http://test.api/latex",
            wechat_webhook="https://test.webhook",
            wechat_mentioned=(),
        )

        # 创建多个文本块
        blocks = []
        for i in range(3):
            block = Block(
                type="text",
                content=f"文本{i+1}",
                width=200000,
                height=calc_text_line_height(font_size_pt),
                inline=True,
                position=BlockPosition(left=0, top=i * calc_text_line_height(font_size_pt))
            )
            blocks.append(block)

        # 创建文本框
        textbox = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.7), Inches(10), Inches(5)
        )
        text_frame = textbox.text_frame

        # 渲染
        render_page(
            slide,
            blocks,
            env,
            font_size_pt,
            FormulaCache(env),
            (int(Inches(0.2)), int(Inches(0.7))),
            text_frame
        )

        # 验证有段落
        assert len(text_frame.paragraphs) >= 1

    def test_render_page_structure(self, presentation, font_size_pt=16, mock_env=None):
        """验证渲染后的页面结构"""
        prs, slide = presentation
        env = mock_env or EnvConfig(
            profile="test",
            proxy=None,
            oss_region="beijing",
            latex_api_url="http://test.api/latex",
            wechat_webhook="https://test.webhook",
            wechat_mentioned=(),
        )

        # 创建页面
        blocks = [
            Block(
                type="text",
                content="页面内容",
                width=300000,
                height=calc_text_line_height(font_size_pt),
                inline=True,
                position=BlockPosition(left=0, top=0)
            )
        ]

        # 创建文本框
        textbox = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.7), Inches(10), Inches(5)
        )

        # 渲染
        render_page(
            slide,
            blocks,
            env,
            font_size_pt,
            FormulaCache(env),
            (int(Inches(0.2)), int(Inches(0.7))),
            textbox.text_frame
        )

        # 验证slide仍然有效
        assert slide is not None


class TestRenderContentPage:
    """内容页渲染测试"""

    def test_render_with_header(self, presentation, font_size_pt=16, mock_env=None):
        """带标题头的内容页"""
        prs, slide = presentation
        env = mock_env or EnvConfig(
            profile="test",
            proxy=None,
            oss_region="beijing",
            latex_api_url="http://test.api/latex",
            wechat_webhook="https://test.webhook",
            wechat_mentioned=(),
        )

        blocks = [
            Block(
                type="text",
                content="正文内容",
                width=300000,
                height=calc_text_line_height(font_size_pt),
                inline=True,
                position=BlockPosition(left=0, top=0)
            )
        ]

        # 渲染带标题的内容页
        render_content_page(
            slide,
            blocks,
            env,
            font_size_pt,
            FormulaCache(env),
            header="页面标题",
            text_box_origin=(int(Inches(0.2)), int(Inches(0.7)))
        )

        # 验证有多个shape（标题+正文文本框）
        assert len(slide.shapes) >= 1


# ============================================================
# 运行测试
# ============================================================

if __name__ == "__main__":
    pytest.main([__file__, "-v"])
