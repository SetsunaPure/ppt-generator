"""
PPT生成器 - PPT构建主流程

负责：
1. 加载PPT模板
2. 按Page列表生成对应的slide
3. 调用渲染引擎填充内容
4. 保存输出文件
"""

import logging
from pathlib import Path
from typing import List, Optional, Dict, Any
from datetime import datetime
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .config import SLIDE_LAYOUTS, PageType, TEXT_BOX_CONFIG, DEFAULT_LOGO_PATH, update_layout_from_template
from .models import Page, PageData, EnvConfig
from .renderer import render_page, render_content_page, FormulaCache
from .utils import (
    get_logger,
    ensure_dir,
    get_timestamp_filename,
    download_image_to_bytes,
    save_temp_image,
)
from .formula_converter import batch_convert_formulas

logger = get_logger()


# ============================================================
# 排版模板生成器
# ============================================================

def generate_layout_template(
    origin_template_path: str,
    output_path: str,
    pages: List[Page],
    logo_url: Optional[str] = None,
    env: Optional[EnvConfig] = None,
    prefix: str = "lesson"
) -> str:
    """
    生成排版模板

    流程：
    1. 解析logo路径（URL则下载到临时文件，本地路径直用）
    2. 计算logo缩放尺寸（固定高度1.07cm，按原始宽高比算宽度）
    3. 遍历pages，按page.type从SLIDE_LAYOUTS取source_slide_index
    4. 复制对应源slide到新Presentation
    5. 每页右上角添加logo
    6. 保存为排版模板文件
    7. 清理临时logo文件
    """
    # 确保输出目录存在
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    # 处理logo
    logo_path = None
    if logo_url:
        logo_path = _download_logo(logo_url, env)
    
    # 如果没有提供logo，使用默认logo
    if not logo_path and Path(DEFAULT_LOGO_PATH).exists():
        logo_path = DEFAULT_LOGO_PATH
        logger.info(f"使用默认logo: {DEFAULT_LOGO_PATH}")

    try:
        # 加载原始模板
        prs = Presentation(origin_template_path)

        # 创建新的Presentation
        new_prs = Presentation()
        new_prs.slide_width = prs.slide_width
        new_prs.slide_height = prs.slide_height

        # 遍历pages生成slides
        for page in pages:
            # 获取source slide index
            layout_config = _get_layout_config(page.type)
            source_idx = layout_config.source_slide_index

            if source_idx < len(prs.slides):
                source_slide = prs.slides[source_idx]
                # 复制slide到新Presentation
                slide = _copy_slide(source_slide, new_prs)

                # 添加logo
                if logo_path:
                    _add_logo_to_slide(slide, logo_path, new_prs.slide_width)
            else:
                # 使用空白布局
                blank_layout = new_prs.slide_layouts[6]  # 空白布局
                slide = new_prs.slides.add_slide(blank_layout)

                # 添加logo
                if logo_path:
                    _add_logo_to_slide(slide, logo_path, new_prs.slide_width)

        # 保存
        new_prs.save(output_path)
        return output_path

    except Exception as e:
        logger.error(f"排版模板生成失败: {e}")
        raise


def _download_logo(logo_url: str, env: Optional[EnvConfig]) -> Optional[str]:
    """下载logo到临时文件"""
    if not env:
        return None

    try:
        image_bytes = download_image_to_bytes(logo_url, env)
        if image_bytes:
            return save_temp_image(image_bytes, ".png")
    except Exception as e:
        logger.warning(f"Logo下载失败: {e}")

    return None


def _get_layout_config(page_type: PageType):
    """获取页面类型的布局配置"""
    type_map = {
        PageType.TITLE: "title",
        PageType.CATALOG: "catalog",
        PageType.FORMULA: "content",
        PageType.QUESTION: "content",
        PageType.ANSWER: "content",
    }
    config_name = type_map.get(page_type, "content")
    return SLIDE_LAYOUTS.get(config_name, SLIDE_LAYOUTS["content"])


def _copy_slide(source_slide, new_prs: Presentation):
    """
    复制slide到新Presentation，保留slide_layout背景和图片关系
    
    核心逻辑（参考老代码duplicate_slide_simple）：
    1. 用源slide的slide_layout创建新slide → 保留master背景
    2. 深拷贝spTree下的shape元素（sp/pic/grpSp/cxnSp）
    3. 修复rId映射（图片等嵌入资源才能正确关联）
    """
    import copy
    from pptx.oxml.ns import qn

    # 1. 用源slide的slide_layout创建新slide（保留背景）
    slide_layout = source_slide.slide_layout
    new_slide = new_prs.slides.add_slide(slide_layout)

    # 1.5 清除layout自带的占位符shape（如Title/Subtitle等）
    # 某些模板的slide_layout自带占位符，会导致拷贝后的shape索引偏移
    target_spTree = new_slide.element.find(
        './/{http://schemas.openxmlformats.org/presentationml/2006/main}spTree'
    )
    if target_spTree is not None:
        # 收集所有需要删除的占位符元素（sp/cxSpGrp等带ph属性的）
        to_remove = []
        for child in list(target_spTree):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            # 只检查sp和grpSp，寻找含placeholder的元素
            if tag in ('sp', 'grpSp', 'cxnSp', 'graphicFrame'):
                # 检查是否包含ph（placeholder）属性
                ph_elements = child.findall(
                    './/{http://schemas.openxmlformats.org/presentationml/2006/main}ph'
                )
                # 也检查nvSpPr/nvGrpSpPr里的ph
                if not ph_elements:
                    ph_elements = child.findall(
                        './/{http://schemas.openxmlformats.org/drawingml/2006/main}ph'
                    )
                if ph_elements:
                    to_remove.append(child)
        for elem in to_remove:
            target_spTree.remove(elem)
        if to_remove:
            logger.debug(f"清除了 {len(to_remove)} 个layout占位符shape")

    # 2. 复制源slide的shape元素到新slide
    source_spTree = source_slide.element.find(
        './/{http://schemas.openxmlformats.org/presentationml/2006/main}spTree'
    )

    if source_spTree is not None and target_spTree is not None:
        copied_count = 0
        for child in source_spTree:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ['sp', 'pic', 'grpSp', 'cxnSp']:
                copied_child = copy.deepcopy(child)
                target_spTree.append(copied_child)
                copied_count += 1
    # 3. 修复rId映射：建立源slide → 新slide的关系映射
    rId_map = {}
    for rId, rel in list(source_slide.part.rels.items()):
        if "notesSlide" in rel.reltype:
            continue
        new_rId = new_slide.part.relate_to(rel._target, rel.reltype, rel.is_external)
        rId_map[rId] = new_rId

    # 替换所有blip中的rId引用
    def fix_rid(element):
        for blip in element.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
            old_rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if old_rId and old_rId in rId_map:
                blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', rId_map[old_rId])

    if target_spTree is not None:
        for child in target_spTree:
            fix_rid(child)

    return new_slide


def _add_logo_to_slide(slide, logo_path: str, slide_width, logo_height_inches: float = 0.42):
    """添加logo到slide右上角"""
    try:
        from PIL import Image as PILImage

        # 获取logo尺寸
        with PILImage.open(logo_path) as img:
            orig_w, orig_h = img.size

        # 按比例计算宽度
        aspect_ratio = orig_w / orig_h if orig_h > 0 else 1.0
        logo_width = logo_height_inches * aspect_ratio

        # 定位：右上角，间距0.2in
        left = slide_width - int(Inches(logo_width + 0.2))
        top = int(Inches(0.2))

        # 添加图片
        slide.shapes.add_picture(
            logo_path,
            left,
            top,
            int(Inches(logo_width)),
            int(Inches(logo_height_inches))
        )
    except Exception as e:
        logger.warning(f"Logo添加失败: {e}")


# ============================================================
# PPT构建器
# ============================================================

class PptBuilder:
    """
    PPT构建器

    负责完整的PPT生成流程
    """

    def __init__(
        self,
        env: EnvConfig,
        pages: List[Page],
        font_size: int,
        file_content_style: str,
        logo_path: str,
        template_prefix: str = "lesson"
    ):
        self.env = env
        self.pages = pages
        self.font_size = font_size
        self.style = file_content_style
        self.logo_path = logo_path
        self.prefix = template_prefix

        # 初始化公式缓存
        self.formula_cache = FormulaCache(env)

        # 确保目录存在
        self._ensure_dirs()

    def _ensure_dirs(self):
        """确保目录存在"""
        ensure_dir("./file/output/")
        ensure_dir("./file/template/")
        ensure_dir("./file/layout_template/")
        ensure_dir("./file/log/")

    def build(self) -> str:
        """
        构建PPT

        返回: 输出文件路径
        """
        # 生成路径
        timestamp = get_timestamp_filename("", "")
        output_filename = f"{self.prefix}_output_{timestamp}.pptx"
        output_path = f"./file/output/{output_filename}"

        layout_filename = f"{self.prefix}_template_layout_{timestamp}.pptx"
        layout_path = f"./file/layout_template/{layout_filename}"

        template_filename = f"{self.prefix}_template_{self.style}.pptx"
        template_path = f"./file/template/{self.prefix}/{template_filename}"

        try:
            # 如果模板文件不存在，创建默认模板
            if not Path(template_path).exists():
                self._create_default_template(template_path)

            # 从模板动态读取布局尺寸
            layout_info = update_layout_from_template(template_path)
            logger.info(f"模板布局: slide={layout_info['slide_width']/914400:.1f}x{layout_info['slide_height']/914400:.1f}in, "
                       f"content_area={layout_info['text_box_width']/914400:.1f}x{layout_info['text_box_height']/914400:.1f}in")

            # 更新TEXT_BOX_CONFIG
            from .config import TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT
            TEXT_BOX_CONFIG["left"] = layout_info["content_left"]
            TEXT_BOX_CONFIG["top"] = layout_info["content_top"]
            TEXT_BOX_CONFIG["width"] = layout_info["text_box_width"]
            TEXT_BOX_CONFIG["height"] = layout_info["text_box_height"]

            # 生成排版模板
            logo_url = self.logo_path if self.logo_path else None
            generate_layout_template(
                origin_template_path=template_path,
                output_path=layout_path,
                pages=self.pages,
                logo_url=logo_url,
                env=self.env,
                prefix=self.prefix
            )

            # 加载排版模板
            prs = Presentation(layout_path)

            # 收集所有公式并批量转换
            self._batch_convert_formulas()

            # 遍历pages渲染内容
            for idx, page in enumerate(self.pages):
                if idx < len(prs.slides):
                    slide = prs.slides[idx]
                    self._fill_slide(slide, page, idx)

            # 保存PPT
            prs.save(output_path)
            logger.info(f"PPT生成成功: {output_path}")

            return output_path

        except Exception as e:
            logger.error(f"PPT构建失败: {e}")
            raise

    def _create_default_template(self, template_path: str):
        """创建默认PPT模板"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 宽屏
        prs.slide_height = Inches(7.5)

        # 标题页
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 添加标题文本框
        title_box = title_slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(12), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "标题"

        # 内容页1
        content_slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        content_box1 = content_slide1.shapes.add_textbox(
            Inches(0.2), Inches(0.5), Inches(12.9), Inches(0.8)
        )
        tf1 = content_box1.text_frame
        p1 = tf1.paragraphs[0]
        p1.text = "页眉"

        # 目录页（占位）
        catalog_slide = prs.slides.add_slide(prs.slide_layouts[6])

        prs.save(template_path)

    def _batch_convert_formulas(self):
        """收集并批量转换公式"""
        all_formulas = set()

        for page in self.pages:
            if page.data.content:
                for block in page.data.content:
                    if block.type == "latex" and block.content:
                        all_formulas.add(block.content)

        if all_formulas:
            results = batch_convert_formulas(list(all_formulas), self.env)
            for latex, omml in results.items():
                self.formula_cache.put(latex, omml)

    def _fill_slide(self, slide, page: Page, slide_idx: int):
        """根据 page.type 填充slide内容"""
        layout_config = _get_layout_config(page.type)

        # 按映射填充shape（支持索引和智能查找两种模式）
        for shape_mapping in layout_config.shapes:
            shape = self._find_shape(slide, shape_mapping, slide_idx)
            if shape:
                self._fill_shape(shape, page.data, shape_mapping.field)

        # 内容页需要额外处理
        if layout_config.auto_create_textbox and page.data.content:
            # 创建正文文本框
            textbox = slide.shapes.add_textbox(
                int(TEXT_BOX_CONFIG["left"]),
                int(TEXT_BOX_CONFIG["top"]),
                int(TEXT_BOX_CONFIG["width"]),
                int(TEXT_BOX_CONFIG["height"])
            )
            text_frame = textbox.text_frame
            text_frame.word_wrap = TEXT_BOX_CONFIG["word_wrap"]

            # 渲染内容块
            text_box_origin = (int(TEXT_BOX_CONFIG["left"]), int(TEXT_BOX_CONFIG["top"]))
            render_page(
                slide,
                page.data.content,
                self.env,
                self.font_size,
                self.formula_cache,
                text_box_origin,
                text_frame
            )

    def _find_shape(self, slide, shape_mapping, slide_idx: int):
        """
        智能查找shape：先按索引，索引不匹配时按field语义查找
        
        不同模板的shape顺序可能不同（如template_4的header在index 3而非2），
        索引找不到时回退到按名称/位置/类型语义匹配
        """
        # 优先按索引查找
        idx = shape_mapping.shape_index
        if idx < len(slide.shapes):
            shape = slide.shapes[idx]
            # 验证shape类型与field匹配
            field = shape_mapping.field
            if self._shape_matches_field(shape, field):
                return shape
            # 索引不匹配，走语义查找
            logger.debug(f"Slide {slide_idx}: shape[{idx}] '{shape.name}' 不匹配 field='{field}'，回退语义查找")

        # 语义查找：按field含义匹配shape
        return self._find_shape_by_semantic(slide, shape_mapping.field, slide_idx)

    def _shape_matches_field(self, shape, field: str) -> bool:
        """判断shape是否符合field的预期"""
        if field == "header":
            # header = 标题栏文字，必须是文本框且在顶部区域
            return (hasattr(shape, 'has_text_frame') and shape.has_text_frame
                    and shape.shape_type == 17  # TEXT_BOX
                    and shape.top < 914400)  # top < 1 inch
        elif field == "title":
            return (hasattr(shape, 'has_text_frame') and shape.has_text_frame
                    and shape.shape_type == 17)
        elif field == "subtitle":
            return (hasattr(shape, 'has_text_frame') and shape.has_text_frame
                    and shape.shape_type == 17)
        elif field == "content":
            return (hasattr(shape, 'has_text_frame') and shape.has_text_frame
                    and shape.shape_type == 17)
        return True

    def _find_shape_by_semantic(self, slide, field: str, slide_idx: int):
        """按field语义在slide中查找匹配的shape"""
        if field == "header":
            # 查找顶部区域(top < 1in)的文本框，且不是圆角矩形(AUTO_SHAPE)
            for shape in slide.shapes:
                if (hasattr(shape, 'has_text_frame') and shape.has_text_frame
                        and shape.shape_type == 17  # TEXT_BOX
                        and shape.top < 914400):  # top < 1 inch
                    return shape
        elif field == "title":
            # 标题页的主标题文本框（第二个文本框，因为第一个可能是单元名称）
            text_boxes = [s for s in slide.shapes
                          if hasattr(s, 'has_text_frame') and s.has_text_frame
                          and s.shape_type == 17]
            if len(text_boxes) >= 1:
                return text_boxes[0]
        elif field == "subtitle":
            text_boxes = [s for s in slide.shapes
                          if hasattr(s, 'has_text_frame') and s.has_text_frame
                          and s.shape_type == 17]
            if len(text_boxes) >= 2:
                return text_boxes[1]
        elif field == "content":
            text_boxes = [s for s in slide.shapes
                          if hasattr(s, 'has_text_frame') and s.has_text_frame
                          and s.shape_type == 17]
            if text_boxes:
                return text_boxes[-1]

        logger.warning(f"Slide {slide_idx}: 未找到匹配 field='{field}' 的shape")
        return None

    def _fill_shape(self, shape, data: PageData, field: str):
        """
        填充shape的内容，保留模板原有格式（字体、字号、颜色、粗体等）
        
        核心逻辑（参考老代码copy_text）：
        1. 深拷贝模板段落的run格式
        2. 清空文本框
        3. 用模板run的格式属性创建新run，只替换文本内容
        """
        import copy
        from pptx.util import Pt

        if field == "title":
            text = data.title if data.title else ""
        elif field == "subtitle":
            text = data.subTitle or ""
        elif field == "header":
            text = data.title if data.title else ""
        elif field == "content":
            # 目录页：显示课时名称
            text = data.title if data.title else ""
        else:
            text = ""

        # 设置文本，保留模板格式
        if hasattr(shape, "text_frame"):
            tf = shape.text_frame
            
            # 深拷贝模板段落的格式
            tmp_paragraphs = copy.deepcopy(tf.paragraphs)
            tf.clear()

            if tmp_paragraphs:
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                p.clear()

                # 用模板第一个段落的run格式创建新run，只替换文本
                first_paragraph = tmp_paragraphs[0]
                for run in first_paragraph.runs:
                    r = p.add_run()
                    r.text = text  # 替换文本

                    # 保留模板格式
                    r.font.name = run.font.name if run.font.name is not None else '黑体'
                    r.font.size = run.font.size if run.font.size is not None else Pt(self.font_size)
                    r.font.bold = run.font.bold
                    r.font.italic = run.font.italic
                    r.font.underline = run.font.underline

                    # 保留颜色
                    if run.font.color.type == 1:  # RGB颜色
                        r.font.color.rgb = run.font.color.rgb
                    elif run.font.color.type == 2:  # 主题颜色（scheme）
                        r.font.color.theme_color = run.font.color.theme_color
                    elif run.font.color.type == 3:  # 主题颜色+色调
                        r.font.color.theme_color = run.font.color.theme_color
                        r.font.color.tint_and_shade = run.font.color.tint_and_shade
            else:
                # 无模板段落，直接设文本
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                r = p.add_run()
                r.text = text
                r.font.name = "黑体"
                r.font.size = Pt(self.font_size)


# ============================================================
# 便捷函数
# ============================================================

def build_ppt(
    pages: List[Page],
    output_path: str,
    env: EnvConfig,
    font_size: int = 16,
    style: str = "0",
    logo_path: str = "",
    prefix: str = "lesson"
) -> str:
    """
    快速构建PPT

    参数:
        pages: Page列表
        output_path: 输出路径
        env: 环境配置
        font_size: 字号
        style: 模板风格
        logo_path: logo路径
        prefix: 前缀(lesson/topic)

    返回: 输出文件路径
    """
    builder = PptBuilder(
        env=env,
        pages=pages,
        font_size=font_size,
        file_content_style=style,
        logo_path=logo_path,
        template_prefix=prefix
    )
    return builder.build()
