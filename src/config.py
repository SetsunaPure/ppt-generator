"""
PPT生成器 - 配置模块

代理环境配置化（生产环境走代理，dev/test直连）
shapes索引映射配置化（换模板只改配置）
"""

from dataclasses import dataclass, field
from typing import Optional, Dict, Any
from enum import Enum


# ============================================================
# 代理配置
# ============================================================

PROXY_SETTINGS = {
    "prod": {
        "http": "http://10.1.1.219:8500",
        "https": "http://10.1.1.219:8500",
    },
    "dev": None,
    "test": None,
}


def get_proxy_config(profile: str) -> Optional[Dict[str, str]]:
    """根据环境标识获取代理配置，未知profile返回None"""
    return PROXY_SETTINGS.get(profile)


# ============================================================
# 布局常量
# ============================================================

from pptx.util import Inches, Pt, Emu

# 注意：以下值为默认值，实际运行时从模板动态读取（见 update_layout_from_template）
TEXT_BOX_WIDTH = Inches(9.6)                 # 文本框宽度（模板10in - 左右各0.2in）
TEXT_BOX_HEIGHT = Inches(4.9)                # 文本框高度（标题栏底部到slide底部）
PADDING_LEFT_RIGHT = Pt(10)                  # 左右内边距（与text_frame.margin对齐）
AVAILABLE_WIDTH = TEXT_BOX_WIDTH - PADDING_LEFT_RIGHT * 2  # 实际可用内容宽度

# 默认logo路径
DEFAULT_LOGO_PATH = "./file/logo/logo_canpoint.png"


def update_layout_from_template(template_path: str) -> dict:
    """
    从PPT模板文件动态读取布局尺寸，更新全局常量
    
    返回: 更新后的布局参数字典
    
    流程:
    1. 读取模板的slide_width/slide_height
    2. 读取正文slide(索引1)的标题栏shape，计算内容区域起始Y
    3. 更新 TEXT_BOX_WIDTH / TEXT_BOX_HEIGHT / AVAILABLE_WIDTH
    """
    from pptx import Presentation
    
    global TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT, AVAILABLE_WIDTH
    
    prs = Presentation(template_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 从正文slide(索引1)读取标题栏位置
    # 优先级：圆角矩形标题栏(AUTO_SHAPE) > 标题文本框 > 默认值
    if len(prs.slides) > 1:
        content_slide = prs.slides[1]
        header_bottom_emu = Inches(0.7)  # 默认值（与老代码一致）
        header_left_emu = Inches(0.2)
        found_header = False
        
        # 第一轮：找圆角矩形标题栏（AUTO_SHAPE，顶部区域，高度<1in）
        for shape in content_slide.shapes:
            if shape.shape_type == 1:  # AUTO_SHAPE
                if shape.top < Inches(1) and shape.height < Inches(1):
                    header_bottom_emu = shape.top + shape.height
                    header_left_emu = shape.left
                    found_header = True
                    break
        
        # 第二轮：fallback到标题文本框
        if not found_header:
            for shape in content_slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt and shape.top < Inches(1):
                        header_bottom_emu = shape.top + shape.height
                        break
    else:
        header_bottom_emu = Inches(0.7)
        header_left_emu = Inches(0.2)
    
    # 计算文本框尺寸
    margin_lr = Inches(0.2)
    TEXT_BOX_WIDTH = slide_width - margin_lr * 2
    TEXT_BOX_HEIGHT = slide_height - header_bottom_emu - Inches(0.1)
    AVAILABLE_WIDTH = TEXT_BOX_WIDTH - PADDING_LEFT_RIGHT * 2
    
    layout_info = {
        'slide_width': slide_width,
        'slide_height': slide_height,
        'text_box_width': TEXT_BOX_WIDTH,
        'text_box_height': TEXT_BOX_HEIGHT,
        'content_top': header_bottom_emu,
        'content_left': margin_lr,
    }
    
    return layout_info

LINE_SPACING_FACTOR = 1.3                    # 行距系数（1.3倍行距）


def calc_text_line_height(font_size_pt: int) -> int:
    """
    根据字号计算单行文字的高度(EMU)

    注意：这只是一个文字行的高度增量，不是"每页几行"。
    在混排中，图片/表格占用的垂直空间与字号无关。
    页满判断由布局引擎统一用 cursor_y vs TEXT_BOX_HEIGHT 决定。
    """
    return int(Pt(font_size_pt) * LINE_SPACING_FACTOR)


# ============================================================
# 图片尺寸级别
# ============================================================

class ImageSizeLevel(str, Enum):
    """图片尺寸级别枚举"""
    NORMAL = "normal"       # 常规图：缩放后高度 < 60% TEXT_BOX_HEIGHT，可跟文字同页
    LARGE = "large"         # 大图：缩放后高度 ≥ 60% TEXT_BOX_HEIGHT，独占一页
    FULLPAGE = "fullpage"   # 超大图：原始高度本身就超过TEXT_BOX_HEIGHT，居中撑满


# ============================================================
# Page类型枚举
# ============================================================

class PageType(str, Enum):
    TITLE = "title"         # 标题页
    CATALOG = "ml"          # 目录页
    FORMULA = "formula"     # 公式内容页
    QUESTION = "question"   # 题目页
    ANSWER = "answer"       # 答案页


# ============================================================
# Slide布局配置
# ============================================================

@dataclass(frozen=True)
class ShapeMapping:
    """单个shape的映射"""
    shape_index: int               # 在slide.shapes中的索引
    field: str                     # 填充字段: title/subtitle/header/content


@dataclass(frozen=True)
class SlideLayout:
    """一种slide类型的shape映射"""
    source_slide_index: int        # 模板中源slide索引
    shapes: tuple[ShapeMapping]    # shape映射列表（使用tuple保证不可变）
    auto_create_textbox: bool = False  # 是否自动创建正文文本框


SLIDE_LAYOUTS: Dict[str, SlideLayout] = {
    "title": SlideLayout(
        source_slide_index=0,
        shapes=(
            ShapeMapping(shape_index=1, field="title"),
            ShapeMapping(shape_index=2, field="subtitle"),
        ),
    ),
    "catalog": SlideLayout(
        source_slide_index=2,
        shapes=(
            ShapeMapping(shape_index=2, field="content"),
        ),
    ),
    "content": SlideLayout(
        source_slide_index=1,
        shapes=(
            ShapeMapping(shape_index=2, field="header"),
        ),
        auto_create_textbox=True,
    ),
}


# ============================================================
# 文本框自动创建参数
# ============================================================

TEXT_BOX_CONFIG = {
    "left": Inches(0.2),
    "top": Inches(0.7),
    "width": Inches(9.6),
    "height": Inches(4.9),
    "margin": Pt(10),
    "font_name": "黑体",
    "word_wrap": True,
}


# ============================================================
# 字符宽度规则
# ============================================================

def char_width_factor(ch: str) -> float:
    """
    按字符宽度规则计算字符宽度系数

    基准：1.0 = font_size_emu，即字号高度。
    实测16pt黑体中文字宽≈字号宽度，factor=1.0；ASCII约为中文一半，factor=0.5。

    | 字符类型           | 宽度系数 |
    | 中文汉字           | 1.0     |
    | 中文标点           | 1.0     |
    | 全角符号           | 1.0     |
    | 中文引号           | 1.0     |
    | 中文间隔号·        | 1.0     |
    | 中文省略号…        | 1.0     |
    | 中文破折号—        | 1.0     |
    | 带圈数字/字母      | 1.33    |
    | 其他字符(ASCII等)  | 0.5     |
    """
    code = ord(ch)
    # 中文汉字: \u4e00-\u9fff（16pt黑体中文字约等于字号宽度，factor≈1.0）
    if 0x4e00 <= code <= 0x9fff:
        return 1.0
    # 中文标点: \u3000-\u303f
    if 0x3000 <= code <= 0x303f:
        return 1.0
    # 全角符号: \uff00-\uffef
    if 0xff00 <= code <= 0xffef:
        return 1.0
    # 中文引号: \u2018-\u201d
    if 0x2018 <= code <= 0x201d:
        return 1.0
    # 带圈数字/字母: \u2460-\u24ff
    if 0x2460 <= code <= 0x24ff:
        return 1.33
    # 中文特殊符号
    special_chars = {'·': 1.0, '…': 1.0, '—': 1.0}
    if ch in special_chars:
        return special_chars[ch]
    return 0.5


# ============================================================
# 文件路径配置
# ============================================================

@dataclass
class FilePaths:
    """文件路径配置"""
    template_dir: str = "./file/template/"
    layout_template_dir: str = "./file/layout_template/"
    output_dir: str = "./file/output/"
    log_dir: str = "./file/log/"
    logo_dir: str = "./file/logo/"
