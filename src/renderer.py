"""
PPT生成器 - PPT渲染引擎

核心模块：接收分页器输出的Block列表，渲染到PPT

设计原则：
- Block驱动渲染，每个Block已携带绝对坐标(left, top)和尺寸(width, height)
- 渲染策略按Block类型分道：
  - 行内元素(text/latex)：写入文本框的段落流
  - 块级元素(image/table)：用绝对坐标直接放置到slide上
- 混排协调：文本框段落流中插入占位空段落，为块级元素预留空间
- 答案组渲染：AnswerGroupBlock内部决定2列/4列布局
"""

import logging
from typing import List, Optional, Tuple, Dict, Any
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml import color
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .config import (
    TEXT_BOX_WIDTH,
    TEXT_BOX_HEIGHT,
    TEXT_BOX_CONFIG,
    SLIDE_LAYOUTS,
    ImageSizeLevel,
)
from .models import Block, Page, EnvConfig, AnswerGroupBlock
from .formula_converter import FormulaCache, batch_convert_formulas
from .utils import (
    get_logger,
    download_image_to_bytes,
    save_temp_image,
)
from .compat import clean_html_for_pptx

logger = get_logger()


# ============================================================
# 渲染入口
# ============================================================

def render_page(
    slide,
    blocks: List[Block],
    env: EnvConfig,
    font_size_pt: int,
    formula_cache: FormulaCache,
    text_box_origin: Tuple[int, int],
    text_frame=None,
):
    """
    渲染一页的所有Block，处理混排协调

    参数:
    - slide: python-pptx的Slide对象
    - blocks: 本页Block列表（已有left/top/width/height，按y坐标排序）
    - env: 环境配置
    - font_size_pt: 字号
    - formula_cache: 公式转换缓存
    - text_box_origin: (left_emu, top_emu) 文本框在slide上的起点坐标
    - text_frame: 已创建的文本框frame（可选）

    流程（按y坐标顺序遍历，保证混排对齐）：
    1. 将blocks按y坐标(top)排序，同y按x排序
    2. 遍历每个block：
       a) text/latex → 写入text_frame当前段落
       b) newline → 添加新段落
       c) image/table →
          - 在text_frame中插入占位空段落（space_before=0, line_spacing=精确高度）
          - 在slide上绝对定位放置该元素
       d) 答案组 → 调用render_answer_group
       e) 记录当前y游标，确保后续行内元素从正确位置开始
    3. 所有块级元素绝对定位坐标 = text_box_origin + block.(left, top)
    """
    if not blocks:
        return

    # 按坐标排序
    sorted_blocks = sorted(blocks, key=lambda b: (b.position.top, b.position.left))

    # 如果没有传入text_frame，需要创建
    if text_frame is None:
        # 创建文本框
        left = text_box_origin[0] + int(TEXT_BOX_CONFIG["left"])
        top = text_box_origin[1] + int(TEXT_BOX_CONFIG["top"])
        width = int(TEXT_BOX_CONFIG["width"])
        height = int(TEXT_BOX_CONFIG["height"])

        textbox = slide.shapes.add_textbox(
            left, top, width, height
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = TEXT_BOX_CONFIG["word_wrap"]

    # 初始化段落
    current_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    current_paragraph.clear()

    # 遍历渲染
    current_y = 0
    for block in sorted_blocks:
        block_abs_top = text_box_origin[1] + block.position.top

        if block.inline:
            # === 行内元素 ===
            if block.type == "text":
                _render_text(current_paragraph, block, font_size_pt)
            elif block.type == "latex":
                _render_formula_in_paragraph(current_paragraph, block, font_size_pt, formula_cache)

        elif block.type == "newline":
            # 换行 → 新段落
            current_paragraph = text_frame.add_paragraph()
            current_paragraph.clear()

        elif block.type == "image":
            # === 图片渲染 ===
            # 插入占位空段落
            placeholder_height = block.height
            _add_placeholder_paragraph(text_frame, placeholder_height)

            # 绝对坐标放置图片
            _render_image(slide, block, env, text_box_origin)

        elif block.type == "table":
            # === 表格渲染 ===
            # 插入占位空段落
            placeholder_height = block.height
            _add_placeholder_paragraph(text_frame, placeholder_height)

            # 绝对坐标放置表格
            _render_table(slide, block, font_size_pt, text_box_origin)

        elif block.type == "answer_group":
            # === 答案组渲染 ===
            _render_answer_group(slide, block, font_size_pt, text_box_origin, text_frame, formula_cache)


# ============================================================
# 占位空段落
# ============================================================

def _add_placeholder_paragraph(text_frame, height_emu: int):
    """
    在text_frame中插入占位空段落，为块级元素预留空间

    实现：
    1. text_frame.add_paragraph() 创建空段落
    2. 设置段落行距为精确高度
    3. paragraph.space_before = Pt(0)
    4. paragraph.space_after = Pt(0)
    5. 段落内容为空（无run），纯占位
    """
    from pptx.oxml.ns import qn

    paragraph = text_frame.add_paragraph()
    paragraph.clear()
    # 设置精确行距
    paragraph.line_spacing = Emu(height_emu)
    # python-pptx 不同版本的行距枚举兼容
    try:
        from pptx.enum.text import MSO_LINE_SPACING
        paragraph.line_spacing_rule = MSO_LINE_SPACING.EXACTLY
    except ImportError:
        # 旧版/新版无此枚举，直接通过XML设置
        pPr = paragraph._p.get_or_add_pPr()
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            from lxml import etree
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = lnSpc.find(qn('a:spcPts'))
        if spcPts is None:
            from lxml import etree
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(height_emu / 12700)))  # EMU to half-points
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    return paragraph


# ============================================================
# 文本渲染
# ============================================================

def _render_text(paragraph, block: Block, font_size_pt: int):
    """
    渲染文本Block到段落
    """
    if not block.content:
        return

    # 清洗HTML
    text = clean_html_for_pptx(block.content)

    # 添加run
    run = paragraph.add_run()
    run.text = text
    run.font.name = block.font_name or "黑体"
    run.font.size = Pt(font_size_pt)
    run.font.bold = block.is_bold


def _render_formula_in_paragraph(
    paragraph,
    block: Block,
    font_size_pt: int,
    formula_cache: FormulaCache
):
    """
    渲染LaTeX公式到段落

    优先使用formula_cache中已转换的OMML，否则将公式文本作为普通文本渲染
    """
    if not block.content:
        return

    latex = block.content

    # 尝试从缓存获取已转换的OMML
    omml = formula_cache.get(latex) if formula_cache else None

    if omml:
        # 有OMML结果，尝试插入
        run = paragraph.add_run()
        run.font.name = "Cambria Math"
        run.font.size = Pt(font_size_pt)

        try:
            _insert_omml_into_run(run, omml)
        except Exception as e:
            logger.warning(f"公式OMML插入失败，降级为文本: {e}")
            # 降级：用可读文本代替
            run.text = _latex_to_readable_text(latex)
    else:
        # 无OMML，将公式转为可读文本
        readable = _latex_to_readable_text(latex)
        run = paragraph.add_run()
        run.text = readable
        run.font.name = "Cambria Math"
        run.font.size = Pt(font_size_pt)
        run.font.italic = True  # 公式用斜体区分


def _insert_omml_into_run(run, omml_str: str):
    """
    将OMML XML插入到run中

    格式：
    <a14:m xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main">
        {omml_str}
    </a14:m>
    """
    # 构建完整OMML
    full_omml = (
        f'<a14:m xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main" '
        f'xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">'
        f'{omml_str}'
        f'</a14:m>'
    )

    # 插入到run的XML中
    run._r.append(etree.fromstring(full_omml))


def _latex_to_readable_text(latex: str) -> str:
    """
    将LaTeX公式转为可读纯文本（OMML不可用时的降级方案）

    示例：
    - C_{60} → C₆₀
    - H_2O → H₂O
    - x^2 → x²
    - \\frac{a}{b} → a/b
    - \\boldsymbol{C_{60}} → C₆₀
    """
    import re
    text = latex
    
    # 去掉 \boldsymbol{} 等修饰
    text = re.sub(r'\\boldsymbol\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\mathbf\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\mathrm\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\text\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\textbf\{([^}]*)\}', r'\1', text)
    
    # \frac{a}{b} → a/b
    text = re.sub(r'\\frac\{([^}]*)\}\{([^}]*)\}', r'\1/\2', text)
    
    # 上标 ^{...} 或 ^x
    def replace_sup(match):
        content = match.group(1) or match.group(2)
        sup_map = str.maketrans('0123456789+-=()ni', '⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾ⁿⁱ')
        return content.translate(sup_map)
    text = re.sub(r'\^{([^}]*)}', replace_sup, text)
    text = re.sub(r'\^([^{\\])', replace_sup, text)
    
    # 下标 _{...} 或 _x
    def replace_sub(match):
        content = match.group(1) or match.group(2)
        sub_map = str.maketrans('0123456789+-=()aeiou', '₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎ₐₑᵢₒᵤ')
        return content.translate(sub_map)
    text = re.sub(r'_\{([^}]*)}', replace_sub, text)
    text = re.sub(r'_([^{\\])', replace_sub, text)
    
    # 去掉剩余的常见LaTeX命令
    text = re.sub(r'\\[a-zA-Z]+', '', text)
    text = re.sub(r'[{}]', '', text)
    
    return text.strip()


# ============================================================
# 图片渲染（绝对坐标定位）
# ============================================================

def _render_image(
    slide,
    block: Block,
    env: EnvConfig,
    text_box_origin: Tuple[int, int]
):
    """
    图片用绝对坐标放置到slide上，不嵌入段落流

    参数:
    - block: image类型的Block（含image_url, image_width_inches, image_height_inches, left, top）
    - text_box_origin: 文本框在slide上的起点(left_emu, top_emu)

    流程：
    1. 通过 make_request 下载图片（自动代理）
    2. 保存为临时文件
    3. 计算slide上的绝对坐标
    4. slide.shapes.add_picture() 放置图片
    """
    if not block.image_url:
        return

    # 下载图片
    image_bytes = download_image_to_bytes(block.image_url, env)
    if not image_bytes:
        logger.warning(f"图片下载失败: {block.image_url}")
        return

    # 保存临时文件
    temp_path = save_temp_image(image_bytes, ".png")
    if not temp_path:
        return

    try:
        # 计算绝对坐标
        absolute_left = text_box_origin[0] + block.position.left
        absolute_top = text_box_origin[1] + block.position.top

        # 添加图片
        slide.shapes.add_picture(
            temp_path,
            absolute_left,
            absolute_top,
            Inches(block.image_width_inches or 1),
            Inches(block.image_height_inches or 1)
        )
    except Exception as e:
        logger.error(f"图片添加失败: {e}")
    finally:
        # 清理临时文件
        try:
            Path(temp_path).unlink()
        except Exception:
            pass


# ============================================================
# 表格渲染（绝对坐标定位）
# ============================================================

def _render_table(
    slide,
    block: Block,
    font_size_pt: int,
    text_box_origin: Tuple[int, int]
):
    """
    表格用绝对坐标放置到slide上

    参数:
    - block: table类型的Block（含table_data, table_rows, table_cols, left, top, width, height）
    """
    if not block.table_data:
        return

    # 计算绝对坐标
    absolute_left = text_box_origin[0] + block.position.left
    absolute_top = text_box_origin[1] + block.position.top

    try:
        # 添加表格
        table_shape = slide.shapes.add_table(
            block.table_rows,
            block.table_cols,
            absolute_left,
            absolute_top,
            block.width,
            block.height
        )

        table = table_shape.table

        # 填充单元格内容
        for row_idx, row_data in enumerate(block.table_data):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < block.table_cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_text)

                    # 设置样式
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.name = "黑体"
                    paragraph.font.size = Pt(font_size_pt)
                    paragraph.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                    # 表头样式（第一行）
                    if row_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = color.RGBColor(221, 221, 221)

    except Exception as e:
        logger.error(f"表格渲染失败: {e}")


# ============================================================
# 答案组渲染（AnswerGroupBlock）
# ============================================================

def _render_answer_group(
    slide,
    block: Block,
    font_size_pt: int,
    text_box_origin: Tuple[int, int],
    text_frame,
    formula_cache: FormulaCache
):
    """
    答案组：内部已确定2列/4列布局，渲染时按列数分配X坐标

    参数:
    - block: answer_group类型的Block（含answer_group字段）
    - answer_group.columns: 2 或 4
    - answer_group.items: 答案项inline Block列表

    流程：
    1. 取 group = block.answer_group
    2. absolute_top = text_box_origin[1] + block.top
    3. 根据 columns 计算每个答案项的 X 坐标
    4. 每个item作为inline text写入对应位置
    """
    group = block.answer_group
    if not group:
        return

    absolute_top = text_box_origin[1] + block.position.top
    col_width = TEXT_BOX_WIDTH // group.columns

    # 在text_frame中为答案组添加占位段落
    # 然后按列布局写入答案文本
    for i, item in enumerate(group.items):
        col = i % group.columns
        row = i // group.columns

        # 计算位置
        item_left = text_box_origin[0] + col * col_width
        item_top = absolute_top + row * group.row_height

        # 在slide上创建文本框放置答案
        if i == 0:
            # 复用或创建答案组的文本框
            answer_box = slide.shapes.add_textbox(
                item_left,
                item_top,
                col_width,
                group.row_height
            )
            tf = answer_box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = item.content
            run.font.name = "黑体"
            run.font.size = Pt(font_size_pt)
        else:
            # 添加下一列的答案
            p = tf.add_paragraph()
            p.clear()
            run = p.add_run()
            # 计算缩进使答案对齐到对应列
            indent = col * col_width // 914400  # 转换为英寸
            run.text = " " * (indent * 10) + item.content
            run.font.name = "黑体"
            run.font.size = Pt(font_size_pt)


# ============================================================
# 页面渲染（简化版）
# ============================================================

def render_content_page(
    slide,
    blocks: List[Block],
    env: EnvConfig,
    font_size_pt: int,
    formula_cache: FormulaCache,
    header: str = "",
    text_box_origin: Tuple[int, int] = (Inches(0.2), Inches(0.7))
):
    """
    渲染内容页的简化入口

    包含标题头和内容块
    """
    # 渲染标题
    if header:
        title_box = slide.shapes.add_textbox(
            text_box_origin[0],
            text_box_origin[1],
            TEXT_BOX_WIDTH,
            Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.name = "黑体"
        run.font.size = Pt(font_size_pt + 2)
        run.font.bold = True

    # 内容区域起始Y
    content_origin = (
        text_box_origin[0],
        text_box_origin[1] + Inches(0.5) if header else text_box_origin[1]
    )

    # 渲染内容块
    render_page(
        slide,
        blocks,
        env,
        font_size_pt,
        formula_cache,
        content_origin
    )
